#!/usr/bin/env python3
import argparse
import json
import logging
import shutil
import subprocess
import tempfile
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt
from utils import get_slide_shapes

# Configure logging
logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")
logger = logging.getLogger(__name__)


# --- Analyze Command ---
def handle_analyze(args):
    try:
        prs = Presentation(args.template)
    except Exception as e:
        logger.error(f"Failed to open template: {e}")
        return

    total_slides = len(prs.slides)
    width_inches = prs.slide_width / 914400
    height_inches = prs.slide_height / 914400
    aspect_ratio = width_inches / height_inches
    aspect_str = (
        "16:9 (widescreen)"
        if abs(aspect_ratio - 16 / 9) < 0.01
        else (
            "4:3 (standard)"
            if abs(aspect_ratio - 4 / 3) < 0.01
            else f"{aspect_ratio:.2f}:1"
        )
    )

    print("\n" + "=" * 60)
    print("TEMPLATE ANALYSIS")
    print("=" * 60)
    print(f"Template: {args.template}")
    print(f"Total Slides: {total_slides}")
    print(f"Slide Range: 0-{total_slides - 1} (0-indexed)")
    print(f'Dimensions: {width_inches:.2f}" x {height_inches:.2f}"')
    print(f"Aspect Ratio: {aspect_str}")
    print("=" * 60 + "\n")

    for idx, slide in enumerate(prs.slides):
        shapes = get_slide_shapes(slide)
        placeholders = [
            str(s.shape.placeholder_format.type).split(".")[-1]
            for s in shapes
            if s.shape.is_placeholder
        ]
        preview = " | ".join(
            [
                (
                    s.shape.text.strip()[:50] + "..."
                    if len(s.shape.text) > 50
                    else s.shape.text.strip()
                )
                for s in shapes
                if s.shape.text.strip()
            ][:2]
        )
        print(
            f"Slide {idx}: {len(shapes)} shapes | Placeholders: {', '.join(set(placeholders)) or 'None'}"
        )
        if preview:
            print(f"  Preview: {preview}")


# --- Rearrange Command ---
def duplicate_slide(pres, index):
    source = pres.slides[index]
    new_slide = pres.slides.add_slide(source.slide_layout)
    for _rel_id, rel in source.part.rels.items():
        if "image" in rel.reltype or "media" in rel.reltype:
            new_slide.part.rels.get_or_add(rel.reltype, rel._target)

    # Copy the source slide's background override (`<p:bg>`) if it has
    # one. `add_slide(layout)` only carries over the layout's default
    # background, so a source slide with a per-slide solid/gradient
    # background (e.g. the dark cover RGB 202125 or the black thank-you
    # RGB 000000) was previously rendered with the layout's white default,
    # which broke the contrast the template was designed around.
    p_ns = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    src_csld = source.element.find(f"{p_ns}cSld")
    src_bg = src_csld.find(f"{p_ns}bg") if src_csld is not None else None
    if src_bg is not None:
        dst_csld = new_slide.element.find(f"{p_ns}cSld")
        if dst_csld is not None:
            existing = dst_csld.find(f"{p_ns}bg")
            if existing is not None:
                dst_csld.remove(existing)
            # `<p:bg>` must be the first child of `<p:cSld>`.
            dst_csld.insert(0, deepcopy(src_bg))

    for shape in new_slide.shapes:
        sp = shape.element
        sp.getparent().remove(sp)
    for shape in source.shapes:
        new_el = deepcopy(shape.element)
        # Robustly insert before extLst if it exists, otherwise append
        # This ensures we don't break the schema order but also don't fail if extLst is missing
        ext_lst = new_slide.shapes._spTree.find(
            "{http://schemas.openxmlformats.org/presentationml/2006/main}extLst"
        )
        if ext_lst is not None:
            ext_lst.addprevious(new_el)
        else:
            new_slide.shapes._spTree.append(new_el)

        blips = new_el.xpath(".//a:blip[@r:embed]")
        for blip in blips:
            r_attr = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
            if blip.get(r_attr) in source.part.rels:
                rel = source.part.rels[blip.get(r_attr)]
                blip.set(
                    r_attr, new_slide.part.rels.get_or_add(rel.reltype, rel._target)
                )
    return new_slide


def handle_rearrange(args):
    sequence = [int(x.strip()) for x in args.sequence.split(",")]
    shutil.copy2(args.template, args.output)
    prs = Presentation(args.output)
    total_template_slides = len(prs.slides)
    for idx in sequence:
        if idx >= total_template_slides:
            raise ValueError(f"Index {idx} out of range")

    for i in sequence:
        duplicate_slide(prs, i)
    sldIdLst = prs.slides._sldIdLst
    for _ in range(total_template_slides):
        prs.part.drop_rel(sldIdLst[0].rId)
        del sldIdLst[0]
    prs.save(args.output)
    logger.info(f"Created {args.output} with {len(sequence)} slides")


# --- Inventory Command ---
def handle_inventory(args):
    prs = Presentation(args.input)
    inventory = {}
    for i, slide in enumerate(prs.slides):
        shapes = get_slide_shapes(slide)
        if not shapes:
            continue
        slide_inventory = {}
        for idx, sd in enumerate(shapes):
            sd.shape_id = f"shape-{idx}"
            if not args.issues_only or sd.has_any_issues:
                slide_inventory[sd.shape_id] = sd.to_dict()
        if slide_inventory:
            inventory[f"slide-{i}"] = slide_inventory
    with open(args.output, "w", encoding="utf-8") as f:
        json.dump(inventory, f, indent=2, ensure_ascii=False)
    logger.info(f"Inventory saved to {args.output}")


# --- Replace Command ---
def apply_para_props(paragraph, data):
    pPr = paragraph._element.get_or_add_pPr()
    for child in list(pPr):
        if any(
            child.tag.endswith(s) for s in ["buChar", "buNone", "buAutoNum", "buFont"]
        ):
            pPr.remove(child)
    if data.get("bullet", False):
        level = data.get("level", 0)
        paragraph.level = level
        sz = data.get("font_size", 18.0)
        pPr.set("marL", str(int((sz * (1.6 + level * 1.6)) * 12700)))
        pPr.set("indent", str(int(-sz * 0.8 * 12700)))
        bu = OxmlElement("a:buChar")
        bu.set("char", "•")
        pPr.append(bu)
    else:
        pPr.set("marL", "0")
        pPr.set("indent", "0")
        pPr.insert(0, OxmlElement("a:buNone"))
    if "alignment" in data:
        amap = {
            "LEFT": PP_ALIGN.LEFT,
            "CENTER": PP_ALIGN.CENTER,
            "RIGHT": PP_ALIGN.RIGHT,
            "JUSTIFY": PP_ALIGN.JUSTIFY,
        }
        paragraph.alignment = amap.get(data["alignment"], PP_ALIGN.LEFT)
    if "space_before" in data:
        paragraph.space_before = Pt(data["space_before"])
    if "space_after" in data:
        paragraph.space_after = Pt(data["space_after"])
    run = paragraph.add_run() if not paragraph.runs else paragraph.runs[0]
    run.text = data.get("text", "")
    f = run.font
    if "bold" in data:
        f.bold = data["bold"]
    if "italic" in data:
        f.italic = data["italic"]
    if "font_size" in data:
        f.size = Pt(data["font_size"])
    if "font_name" in data:
        f.name = data["font_name"]
    if "color" in data:
        c = data["color"].lstrip("#")
        if len(c) == 6:
            f.color.rgb = RGBColor(int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))
    elif "theme_color" in data:
        try:
            f.color.theme_color = getattr(MSO_THEME_COLOR, data["theme_color"])
        except AttributeError:
            logger.warning(f"Unknown theme color: {data['theme_color']}")


def _capture_default_style(text_frame) -> dict:
    """Read the first non-empty run's formatting so we can restore it
    after `tf.clear()`. Returns a partial replacement-style dict
    (keys: color, theme_color, font_name, font_size, bold, italic).
    """
    style: dict = {}
    for para in text_frame.paragraphs:
        for run in para.runs:
            f = run.font
            if "font_size" not in style and f.size is not None:
                style["font_size"] = f.size.pt
            if "font_name" not in style and f.name:
                style["font_name"] = f.name
            if "bold" not in style and f.bold is not None:
                style["bold"] = f.bold
            if "italic" not in style and f.italic is not None:
                style["italic"] = f.italic
            if "color" not in style and "theme_color" not in style:
                try:
                    rgb = f.color.rgb
                    if rgb is not None:
                        style["color"] = str(rgb)
                except (AttributeError, TypeError):
                    try:
                        tc = f.color.theme_color
                        if tc is not None:
                            style["theme_color"] = tc.name
                    except (AttributeError, TypeError):
                        pass
            if "color" in style or "theme_color" in style or "font_size" in style:
                return style
    return style


def _slide_default_text_color(slide) -> str | None:
    """Pick a sensible default text color for empty placeholders based on
    the slide's background luminance. Returns a hex color or None when
    the background is light enough that the default (black) works.
    """
    try:
        fill = slide.background.fill
        rgb = None
        try:
            rgb = fill.fore_color.rgb
        except (AttributeError, TypeError):
            rgb = None
        # Fall back to layout master if the slide itself inherits.
        if rgb is None:
            try:
                rgb = slide.slide_layout.background.fill.fore_color.rgb
            except (AttributeError, TypeError):
                rgb = None
        if rgb is None:
            return None
        # Relative luminance (sRGB approximation).
        r, g, b = (int(c) for c in (rgb[0], rgb[1], rgb[2]))
        luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255.0
        if luminance < 0.5:
            return "FFFFFF"
        return None
    except Exception:
        return None


def handle_replace(args):
    prs = Presentation(args.input)
    with open(args.replacements) as f:
        replacements = json.load(f)

    # Validate every reference up front. Silent skips were the historical
    # foot-gun: a wrong slide-N / shape-N would let the LLM ship a deck
    # of "Lorem ipsum" template text under the impression replacement
    # had succeeded.
    errors: list[str] = []
    if not replacements:
        errors.append("replacements dict is empty — nothing to apply")
    for slide_key, shape_data in replacements.items():
        try:
            s_idx = int(slide_key.split("-")[1])
        except (IndexError, ValueError):
            errors.append(f"invalid slide key {slide_key!r} (expected 'slide-<N>')")
            continue
        if s_idx >= len(prs.slides):
            errors.append(
                f"slide-{s_idx} out of range (deck has {len(prs.slides)} slides)"
            )
            continue
        shapes = get_slide_shapes(prs.slides[s_idx])
        for shape_key, r_data in shape_data.items():
            try:
                sh_idx = int(shape_key.split("-")[1])
            except (IndexError, ValueError):
                errors.append(
                    f"invalid shape key {shape_key!r} on slide-{s_idx}"
                    " (expected 'shape-<N>')"
                )
                continue
            if sh_idx >= len(shapes):
                errors.append(
                    f"slide-{s_idx} / {shape_key} out of range — slide has"
                    f" {len(shapes)} shape(s) (indices 0..{len(shapes) - 1})"
                )
                continue
            paragraphs = r_data.get("paragraphs", [])
            if not paragraphs:
                errors.append(
                    f"slide-{s_idx} / {shape_key} has no paragraphs — would"
                    " blank the shape. Omit the entry to preserve the"
                    " template text, or supply at least one paragraph."
                )

    if errors:
        msg = "Invalid replacements:\n  - " + "\n  - ".join(errors)
        logger.error(msg)
        raise ValueError(msg)

    # Apply replacements and track coverage.
    inventory_shape_count = 0
    for slide in prs.slides:
        inventory_shape_count += len(get_slide_shapes(slide))
    applied_paragraphs = 0
    touched_shapes = 0
    touched_slides = set()
    autosize_failures = 0
    for slide_key, shape_data in replacements.items():
        s_idx = int(slide_key.split("-")[1])
        slide = prs.slides[s_idx]
        shapes = get_slide_shapes(slide)
        slide_default_color = _slide_default_text_color(slide)
        for shape_key, r_data in shape_data.items():
            sh_idx = int(shape_key.split("-")[1])
            tf = shapes[sh_idx].shape.text_frame
            # Capture the original run's formatting BEFORE clearing — we
            # use it as a fallback for properties the replacement payload
            # leaves unspecified. Without this, `tf.clear()` followed by
            # `add_run()` produces a fresh black-default run, which renders
            # as black-on-black on dark cover / thank-you slides whose
            # template runs were originally white.
            preserved = _capture_default_style(tf)
            if "color" not in preserved and "theme_color" not in preserved:
                # Empty placeholder on a dark slide → fall back to the
                # slide-background-derived default (white on dark, black
                # on light) so newly populated body shapes stay legible.
                if slide_default_color:
                    preserved["color"] = slide_default_color
            tf.clear()
            paragraphs = r_data.get("paragraphs", [])
            for i, p_data in enumerate(paragraphs):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                # Replacement data wins; preserved style fills in gaps.
                merged = {**preserved, **p_data}
                apply_para_props(p, merged)
                applied_paragraphs += 1
            # Make replaced shapes auto-shrink text to fit. Without this,
            # any string longer than the placeholder's original capacity
            # spills outside the frame and the deck looks broken.
            try:
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except Exception:
                autosize_failures += 1
            touched_shapes += 1
            touched_slides.add(s_idx)

    prs.save(args.output)
    skipped_shapes = inventory_shape_count - touched_shapes
    summary = (
        f"Applied {applied_paragraphs} paragraph(s) across {touched_shapes}"
        f" shape(s) on {len(touched_slides)} slide(s). {skipped_shapes} shape(s)"
        f" left untouched — they retain their original template text."
        f" Auto-fit (text_to_fit_shape) enabled on replaced shapes so longer"
        f" strings shrink to stay inside the frame."
    )
    logger.info(f"Saved: {args.output}")
    logger.info(summary)
    if autosize_failures:
        logger.warning(
            "auto-fit could not be applied on %d shape(s); long text in those"
            " shapes may still overflow.",
            autosize_failures,
        )
    if skipped_shapes > 0:
        logger.warning(
            "%d shape(s) still hold template placeholder text. If the deck"
            " should NOT show 'Lorem ipsum' or other template content, add a"
            " replacement entry for every shape in the inventory.",
            skipped_shapes,
        )
    # Overflow check
    prs_out = Presentation(args.output)
    overflows = []
    for i, slide in enumerate(prs_out.slides):
        for idx, sd in enumerate(get_slide_shapes(slide)):
            if sd.frame_overflow_bottom:
                overflows.append(
                    f'Slide {i} / Shape {idx}: +{sd.frame_overflow_bottom}"'
                )
    if overflows:
        logger.warning(f"Found {len(overflows)} overflow(s):")
        for m in overflows:
            logger.warning(f"  - {m}")
    if args.cleanup:
        for f in [args.input, args.replacements, "text-inventory.json"]:
            if Path(f).exists():
                Path(f).unlink()
                logger.info(f"Removed: {f}")


# --- Thumbnail Command ---
def handle_thumbnail(args):
    if not shutil.which("soffice") or not shutil.which("pdftoppm"):
        logger.error("Missing libreoffice or poppler-utils")
        return
    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        subprocess.run(
            [
                "soffice",
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(tmp_path),
                args.input,
            ],
            check=True,
            capture_output=True,
        )
        pdf_path = tmp_path / f"{Path(args.input).stem}.pdf"
        subprocess.run(
            ["pdftoppm", "-jpeg", "-r", "100", str(pdf_path), str(tmp_path / "slide")],
            check=True,
            capture_output=True,
        )
        imgs = sorted(tmp_path.glob("slide-*.jpg"))
        from PIL import Image, ImageDraw

        chunk_sz = args.cols * (args.cols + 1)
        for i, start in enumerate(range(0, len(imgs), chunk_sz)):
            chunk = imgs[start : start + chunk_sz]
            with Image.open(chunk[0]) as first:
                thumb_h = int(300 * (first.height / first.width))
            rows = (len(chunk) + args.cols - 1) // args.cols
            grid = Image.new(
                "RGB", (args.cols * 320 + 20, rows * (thumb_h + 40) + 20), "white"
            )
            draw = ImageDraw.Draw(grid)
            for j, p in enumerate(chunk):
                r, c = j // args.cols, j % args.cols
                x, y = c * 320 + 20, r * (thumb_h + 40) + 20
                draw.text((x + 5, y), str(start + j), fill="black")
                with Image.open(p) as im:
                    im.thumbnail((300, thumb_h))
                    grid.paste(im, (x, y + 25))
            out_name = (
                f"{args.prefix}-{i + 1}.jpg"
                if len(imgs) > chunk_sz
                else f"{args.prefix}.jpg"
            )
            grid.save(out_name, quality=95)
            logger.info(f"Created: {out_name}")


# --- Cleanup Command ---
def handle_cleanup(args):
    patterns = [
        "working.pptx",
        "text-inventory.json",
        "replacement-text.json",
        "template-content.md",
        "template-analysis.md",
    ]
    found = [p for p in patterns if Path(p).exists()]
    for d in ["template-thumbnails"]:
        if Path(d).is_dir():
            found.append(d)
    if not found:
        logger.info("No temporary files found")
        return
    if not args.yes:
        print(f"Delete: {', '.join(found)}?")
        resp = input("[y/N]: ").lower()
        if resp not in ["y", "yes"]:
            return
    for p in found:
        if Path(p).is_file():
            Path(p).unlink()
        else:
            shutil.rmtree(p)
        logger.info(f"Removed: {p}")


def main():
    parser = argparse.ArgumentParser(description="Presentation Skill CLI")
    subparsers = parser.add_subparsers(dest="command", required=True)

    p_analyze = subparsers.add_parser("analyze", help="Analyze template")
    p_analyze.add_argument("template", help="Template PPTX")

    p_rearrange = subparsers.add_parser("rearrange", help="Rearrange slides")
    p_rearrange.add_argument("template", help="Template PPTX")
    p_rearrange.add_argument("output", help="Output PPTX")
    p_rearrange.add_argument("sequence", help="Comma-separated indices")

    p_inventory = subparsers.add_parser("inventory", help="Extract inventory")
    p_inventory.add_argument("input", help="Input PPTX")
    p_inventory.add_argument("output", help="Output JSON")
    p_inventory.add_argument("--issues-only", action="store_true")

    p_replace = subparsers.add_parser("replace", help="Apply replacements")
    p_replace.add_argument("input", help="Input PPTX")
    p_replace.add_argument("replacements", help="Replacements JSON")
    p_replace.add_argument("output", help="Output PPTX")
    p_replace.add_argument("--cleanup", action="store_true")

    p_thumb = subparsers.add_parser("thumbnail", help="Create thumbnails")
    p_thumb.add_argument("input", help="Input PPTX")
    p_thumb.add_argument("prefix", nargs="?", default="thumbnails")
    p_thumb.add_argument("--cols", type=int, default=5)

    p_clean = subparsers.add_parser("cleanup", help="Clean up files")
    p_clean.add_argument("--yes", action="store_true")

    args = parser.parse_args()
    cmds = {
        "analyze": handle_analyze,
        "rearrange": handle_rearrange,
        "inventory": handle_inventory,
        "replace": handle_replace,
        "thumbnail": handle_thumbnail,
        "cleanup": handle_cleanup,
    }
    cmds[args.command](args)


if __name__ == "__main__":
    main()
