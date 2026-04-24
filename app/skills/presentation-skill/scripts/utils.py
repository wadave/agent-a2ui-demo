#!/usr/bin/env python3
import logging
import platform
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont
from pptx.enum.text import PP_ALIGN
from pptx.shapes.base import BaseShape

# Configure logging
logger = logging.getLogger(__name__)

# Type aliases
JsonValue = str | int | float | bool | None

ParagraphDict = dict[str, JsonValue]
ShapeDict = dict[str, Any]


@dataclass
class ShapeWithPosition:
    """A shape with its absolute position on the slide."""

    shape: BaseShape
    absolute_left: int  # in EMUs
    absolute_top: int  # in EMUs


class ParagraphData:
    """Data structure for paragraph properties extracted from a PowerPoint paragraph."""

    def __init__(self, paragraph: Any):
        """Initialize from a PowerPoint paragraph object."""
        self.text: str = paragraph.text.strip()
        self.bullet: bool = False
        self.level: int | None = None
        self.alignment: str | None = None
        self.space_before: float | None = None
        self.space_after: float | None = None
        self.font_name: str | None = None
        self.font_size: float | None = None
        self.bold: bool | None = None
        self.italic: bool | None = None
        self.underline: bool | None = None
        self.color: str | None = None
        self.theme_color: str | None = None
        self.line_spacing: float | None = None

        self._extract_properties(paragraph)

    def _extract_properties(self, paragraph: Any):
        """Extract properties from the paragraph object."""
        # Check for bullet formatting
        if hasattr(paragraph, "_p") and paragraph._p.pPr is not None:
            pPr = paragraph._p.pPr
            ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
            if (
                pPr.find(f"{ns}buChar") is not None
                or pPr.find(f"{ns}buAutoNum") is not None
            ):
                self.bullet = True
                if hasattr(paragraph, "level"):
                    self.level = paragraph.level

        # Alignment
        if hasattr(paragraph, "alignment") and paragraph.alignment is not None:
            alignment_map = {
                PP_ALIGN.CENTER: "CENTER",
                PP_ALIGN.RIGHT: "RIGHT",
                PP_ALIGN.JUSTIFY: "JUSTIFY",
            }
            if paragraph.alignment in alignment_map:
                self.alignment = alignment_map[paragraph.alignment]

        # Spacing
        if hasattr(paragraph, "space_before") and paragraph.space_before:
            self.space_before = paragraph.space_before.pt
        if hasattr(paragraph, "space_after") and paragraph.space_after:
            self.space_after = paragraph.space_after.pt

        # Font properties from first run
        if paragraph.runs:
            run = paragraph.runs[0]
            if hasattr(run, "font"):
                font = run.font
                self.font_name = font.name
                if font.size:
                    self.font_size = font.size.pt
                self.bold = font.bold
                self.italic = font.italic
                self.underline = font.underline

                # Color
                try:
                    if font.color.rgb:
                        self.color = str(font.color.rgb)
                except (AttributeError, TypeError):
                    try:
                        if font.color.theme_color:
                            self.theme_color = font.color.theme_color.name
                    except (AttributeError, TypeError):
                        pass

        # Line spacing
        if hasattr(paragraph, "line_spacing") and paragraph.line_spacing is not None:
            if hasattr(paragraph.line_spacing, "pt"):
                self.line_spacing = round(paragraph.line_spacing.pt, 2)
            else:
                base_size = self.font_size or 12.0
                self.line_spacing = round(paragraph.line_spacing * base_size, 2)

    def to_dict(self) -> ParagraphDict:
        """Convert to dictionary for JSON serialization."""
        return {
            k: v
            for k, v in self.__dict__.items()
            if v is not None and not k.startswith("_")
        }


class ShapeData:
    """Data structure for shape properties extracted from a PowerPoint shape."""

    @staticmethod
    def emu_to_inches(emu: int) -> float:
        return emu / 914400.0

    @staticmethod
    def inches_to_pixels(inches: float, dpi: int = 96) -> int:
        return int(inches * dpi)

    @staticmethod
    def get_font_path(font_name: str) -> str | None:
        """Find font file path for a given font name."""
        system = platform.system()
        variants = [
            font_name,
            font_name.lower(),
            font_name.replace(" ", ""),
            font_name.replace(" ", "-"),
        ]

        if system == "Darwin":
            dirs = ["/System/Library/Fonts/", "/Library/Fonts/", "~/Library/Fonts/"]
            exts = [".ttf", ".otf", ".ttc", ".dfont"]
        else:
            dirs = [
                "/usr/share/fonts/truetype/",
                "/usr/local/share/fonts/",
                "~/.fonts/",
            ]
            exts = [".ttf", ".otf"]

        for d in dirs:
            dir_path = Path(d).expanduser()
            if not dir_path.exists():
                continue
            for v in variants:
                for e in exts:
                    p = dir_path / f"{v}{e}"
                    if p.exists():
                        return str(p)

            try:
                for f in dir_path.rglob("*"):
                    if f.is_file() and any(f.name.lower().endswith(e) for e in exts):
                        if font_name.lower().replace(" ", "") in f.name.lower().replace(
                            " ", ""
                        ):
                            return str(f)
            except (OSError, PermissionError):
                continue
        return None

    def __init__(
        self, shape: BaseShape, absolute_left: int, absolute_top: int, slide: Any
    ):
        self.shape = shape
        self.shape_id: str = ""
        self.slide = slide

        # Dimensions
        self.left_emu = absolute_left
        self.top_emu = absolute_top
        self.width_emu = shape.width
        self.height_emu = shape.height

        self.left = round(self.emu_to_inches(self.left_emu), 2)
        self.top = round(self.emu_to_inches(self.top_emu), 2)
        self.width = round(self.emu_to_inches(self.width_emu), 2)
        self.height = round(self.emu_to_inches(self.height_emu), 2)

        # Slide Context
        self.slide_width_emu = (
            slide.part.package.presentation_part.presentation.slide_width
        )
        self.slide_height_emu = (
            slide.part.package.presentation_part.presentation.slide_height
        )

        # Properties
        self.placeholder_type: str | None = None
        self._extract_placeholder_info()

        # Issue Detection
        self.frame_overflow_bottom: float | None = None
        self.slide_overflow_right: float | None = None
        self.slide_overflow_bottom: float | None = None
        self.overlapping_shapes: dict[str, float] = {}
        self.warnings: list[str] = []

        self._analyze_issues()

    def _extract_placeholder_info(self):
        """Extract placeholder type and default font size."""
        if hasattr(self.shape, "is_placeholder") and self.shape.is_placeholder:
            fmt = self.shape.placeholder_format
            if fmt and fmt.type:
                self.placeholder_type = str(fmt.type).split(".")[-1]

    def _analyze_issues(self):
        """Perform issue detection for the shape."""
        self._estimate_frame_overflow()
        self._check_slide_overflow()
        self._check_bullets()

    def _get_default_font_size(self) -> int:
        """Get default font size from theme text styles or use conservative default."""
        try:
            if not (
                hasattr(self.shape, "part") and hasattr(self.shape.part, "slide_layout")
            ):
                return 14

            slide_master = self.shape.part.slide_layout.slide_master
            if not hasattr(slide_master, "element"):
                return 14

            style_name = "bodyStyle"
            if self.placeholder_type and "TITLE" in self.placeholder_type:
                style_name = "titleStyle"

            for child in slide_master.element.iter():
                tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                if tag == style_name:
                    for elem in child.iter():
                        if "sz" in elem.attrib:
                            return int(elem.attrib["sz"]) // 100
        except Exception:
            pass
        return 14

    def _get_usable_dimensions(self, text_frame) -> tuple[int, int]:
        margins = {"top": 0.05, "bottom": 0.05, "left": 0.1, "right": 0.1}
        if hasattr(text_frame, "margin_top") and text_frame.margin_top:
            margins["top"] = self.emu_to_inches(text_frame.margin_top)
        if hasattr(text_frame, "margin_bottom") and text_frame.margin_bottom:
            margins["bottom"] = self.emu_to_inches(text_frame.margin_bottom)
        if hasattr(text_frame, "margin_left") and text_frame.margin_left:
            margins["left"] = self.emu_to_inches(text_frame.margin_left)
        if hasattr(text_frame, "margin_right") and text_frame.margin_right:
            margins["right"] = self.emu_to_inches(text_frame.margin_right)

        usable_width = self.width - margins["left"] - margins["right"]
        usable_height = self.height - margins["top"] - margins["bottom"]

        return (
            self.inches_to_pixels(usable_width),
            self.inches_to_pixels(usable_height),
        )

    def _wrap_text_line(self, line: str, max_width_px: int, draw, font) -> list[str]:
        if not line:
            return [""]
        if draw.textlength(line, font=font) <= max_width_px:
            return [line]
        wrapped = []
        words = line.split(" ")
        current_line = ""
        for word in words:
            test_line = current_line + (" " if current_line else "") + word
            if draw.textlength(test_line, font=font) <= max_width_px:
                current_line = test_line
            else:
                if current_line:
                    wrapped.append(current_line)
                current_line = word
        if current_line:
            wrapped.append(current_line)
        return wrapped

    def _estimate_frame_overflow(self):
        """Estimate text overflow using PIL."""
        if not self.shape or not hasattr(self.shape, "text_frame"):
            return
        text_frame = self.shape.text_frame
        if not text_frame or not text_frame.paragraphs:
            return

        usable_width_px, usable_height_px = self._get_usable_dimensions(text_frame)
        if usable_width_px <= 0 or usable_height_px <= 0:
            return

        dummy_img = Image.new("RGB", (1, 1))
        draw = ImageDraw.Draw(dummy_img)
        default_font_size = self._get_default_font_size()
        total_height_px = 0

        for para_idx, paragraph in enumerate(text_frame.paragraphs):
            if not paragraph.text.strip():
                continue
            para_data = ParagraphData(paragraph)
            font_name = para_data.font_name or "Arial"
            font_size = int(para_data.font_size or default_font_size)
            font_path = self.get_font_path(font_name)
            try:
                font = (
                    ImageFont.truetype(font_path, size=font_size)
                    if font_path
                    else ImageFont.load_default()
                )
            except Exception:
                font = ImageFont.load_default()

            all_wrapped_lines = []
            for line in paragraph.text.split("\n"):
                all_wrapped_lines.extend(
                    self._wrap_text_line(line, usable_width_px, draw, font)
                )

            if all_wrapped_lines:
                line_height_px = (para_data.line_spacing or font_size) * 96 / 72
                if para_idx > 0 and para_data.space_before:
                    total_height_px += para_data.space_before * 96 / 72
                total_height_px += len(all_wrapped_lines) * line_height_px
                if para_data.space_after:
                    total_height_px += para_data.space_after * 96 / 72

        if total_height_px > usable_height_px:
            overflow_inches = round((total_height_px - usable_height_px) / 96.0, 2)
            if overflow_inches > 0.05:
                self.frame_overflow_bottom = overflow_inches

    def _check_slide_overflow(self):
        """Check if shape is outside slide boundaries."""
        if self.left_emu + self.width_emu > self.slide_width_emu:
            self.slide_overflow_right = round(
                self.emu_to_inches(
                    self.left_emu + self.width_emu - self.slide_width_emu
                ),
                2,
            )
        if self.top_emu + self.height_emu > self.slide_height_emu:
            self.slide_overflow_bottom = round(
                self.emu_to_inches(
                    self.top_emu + self.height_emu - self.slide_height_emu
                ),
                2,
            )

    def _check_bullets(self):
        """Detect manual bullet usage."""
        bullet_symbols = ["•", "●", "○"]
        for p in self.shape.text_frame.paragraphs:
            t = p.text.strip()
            if t and any(t.startswith(s + " ") for s in bullet_symbols):
                self.warnings.append(
                    "manual_bullet_symbol: use proper bullet formatting"
                )
                break

    @property
    def paragraphs(self) -> list[ParagraphData]:
        if not hasattr(self.shape, "text_frame"):
            return []
        return [
            ParagraphData(p) for p in self.shape.text_frame.paragraphs if p.text.strip()
        ]

    @property
    def has_any_issues(self) -> bool:
        return any(
            [
                self.frame_overflow_bottom,
                self.slide_overflow_right,
                self.slide_overflow_bottom,
                self.overlapping_shapes,
                self.warnings,
            ]
        )

    def to_dict(self) -> ShapeDict:
        res = {
            "left": self.left,
            "top": self.top,
            "width": self.width,
            "height": self.height,
            "paragraphs": [p.to_dict() for p in self.paragraphs],
        }
        if self.placeholder_type:
            res["placeholder_type"] = self.placeholder_type
        issues = {}
        if self.frame_overflow_bottom:
            issues["frame"] = {"overflow_bottom": self.frame_overflow_bottom}
        slide_off = {
            k: v
            for k, v in {
                "right": self.slide_overflow_right,
                "bottom": self.slide_overflow_bottom,
            }.items()
            if v
        }
        if slide_off:
            issues["slide"] = slide_off
        if self.overlapping_shapes:
            issues["overlap"] = self.overlapping_shapes
        if issues:
            res["issues"] = issues
        if self.warnings:
            res["warnings"] = self.warnings
        return res


def collect_shapes(
    shape: BaseShape, left: int, top: int, slide: Any
) -> list[ShapeData]:
    """Recursively collect text shapes."""
    if hasattr(shape, "shapes"):  # GroupShape
        res = []
        for s in shape.shapes:
            res.extend(collect_shapes(s, left + shape.left, top + shape.top, slide))
        return res

    # Include every shape with a text_frame, even if its current text is
    # empty. Empty placeholders (e.g. the body area on the BLANK_1 content
    # slide, or the entry titles next to each "01"-"05" on the TOC slide)
    # are precisely what the LLM needs to fill — silently dropping them
    # left whole body slides un-addressable and shipped blank.
    if hasattr(shape, "text_frame"):
        if hasattr(shape, "is_placeholder") and shape.is_placeholder:
            if str(shape.placeholder_format.type).endswith("SLIDE_NUMBER"):
                return []
        return [ShapeData(shape, left + shape.left, top + shape.top, slide)]
    return []


def get_slide_shapes(slide: Any) -> list[ShapeData]:
    """Collect and sort shapes from a slide."""
    shapes = []
    for s in slide.shapes:
        shapes.extend(collect_shapes(s, 0, 0, slide))
    shapes.sort(key=lambda s: (s.top // 0.5, s.left))
    return shapes
